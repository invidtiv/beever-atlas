"""Tests for the ingestion pipeline hardening features.

Covers:
- Coreference resolution (Group 1)
- Semantic entity deduplication (Group 2)
- Multimodal expansion (Group 3)
- Semantic search (Group 4)
- Temporal fact lifecycle (Group 5)
- Cross-batch thread context (Group 6)
- Soft orphan handling (Group 7)
"""

from __future__ import annotations

import io
from datetime import UTC, datetime
from unittest.mock import AsyncMock, MagicMock, patch

import pytest

# ── Group 1: Coreference Resolution ──────────────────────────────────


class TestCoreferenceResolver:
    """Tests for coreference resolution service."""

    def test_has_resolvable_references_detects_pronouns(self):
        from beever_atlas.services.coreference_resolver import has_resolvable_references

        messages = [{"text": "Alice built Atlas. It uses Redis."}]
        assert has_resolvable_references(messages) is True

    def test_has_resolvable_references_no_pronouns(self):
        from beever_atlas.services.coreference_resolver import has_resolvable_references

        messages = [{"text": "PostgreSQL version 15 released."}]
        assert has_resolvable_references(messages) is False

    def test_has_resolvable_references_implicit_reference(self):
        from beever_atlas.services.coreference_resolver import has_resolvable_references

        messages = [{"text": "The project needs more resources."}]
        assert has_resolvable_references(messages) is True

    @pytest.mark.asyncio
    async def test_resolve_coreferences_no_pronouns_skips_llm(self):
        from beever_atlas.services.coreference_resolver import resolve_coreferences

        messages = [{"text": "PostgreSQL version 15 released."}]
        result = await resolve_coreferences(messages)
        assert result[0]["raw_text"] == "PostgreSQL version 15 released."
        assert result[0]["text"] == "PostgreSQL version 15 released."

    @pytest.mark.asyncio
    async def test_resolve_coreferences_preserves_raw_text(self):
        """Even when LLM fails, raw_text should be preserved."""
        from beever_atlas.services.coreference_resolver import resolve_coreferences

        messages = [{"text": "Alice built Atlas. It uses Redis."}]
        # Patch the genai client to raise — test graceful fallback
        with patch("beever_atlas.services.coreference_resolver.get_settings") as mock_settings:
            mock_settings.return_value = MagicMock(
                coref_model="gemini-2.5-flash",
                google_api_key="fake",
            )
            with patch("google.genai.Client", side_effect=ImportError("no genai")):
                result = await resolve_coreferences(messages)
        assert result[0]["raw_text"] == "Alice built Atlas. It uses Redis."

    @pytest.mark.asyncio
    async def test_resolve_coreferences_empty_batch(self):
        from beever_atlas.services.coreference_resolver import resolve_coreferences

        result = await resolve_coreferences([])
        assert result == []


# ── Group 2: Semantic Entity Deduplication ────────────────────────────


class TestSemanticEntityDedup:
    """Tests for embedding-based entity deduplication."""

    @pytest.mark.asyncio
    async def test_find_similar_by_embedding(self):
        from beever_atlas.stores.entity_registry import EntityRegistry

        mock_neo4j = AsyncMock()
        mock_neo4j.execute_query = AsyncMock(return_value=[
            {"name": "Beever Atlas", "vec": [0.9, 0.1, 0.0]},
            {"name": "Redis Cache", "vec": [0.1, 0.9, 0.0]},
        ])

        registry = EntityRegistry(mock_neo4j)
        results = await registry.find_similar_by_embedding(
            name="Atlas",
            name_vector=[0.85, 0.15, 0.0],
            threshold=0.8,
        )
        # "Beever Atlas" should be similar, "Redis Cache" should not
        names = [r[0] for r in results]
        assert "Beever Atlas" in names
        assert "Redis Cache" not in names

    def test_merge_rejection_cache(self):
        from beever_atlas.stores.entity_registry import EntityRegistry

        registry = EntityRegistry(MagicMock())
        assert registry.is_merge_rejected("Atlas", "Beever Atlas") is False

        registry.cache_merge_rejection("Atlas", "Beever Atlas")
        assert registry.is_merge_rejected("Atlas", "Beever Atlas") is True
        # Order-independent
        assert registry.is_merge_rejected("Beever Atlas", "Atlas") is True


# ── Group 3: Multimodal Expansion ─────────────────────────────────────


class TestMediaExtractors:
    """Tests for the media extractor registry."""

    def test_registry_dispatches_by_extension(self):
        from beever_atlas.services.media_extractors import (
            OfficeExtractor,
            PdfExtractor,
            create_default_registry,
        )

        registry = create_default_registry()
        assert isinstance(registry.get_extractor(filename="report.pdf"), PdfExtractor)
        assert isinstance(registry.get_extractor(filename="doc.docx"), OfficeExtractor)
        assert isinstance(registry.get_extractor(filename="data.xlsx"), OfficeExtractor)
        assert isinstance(registry.get_extractor(filename="slides.pptx"), OfficeExtractor)

    def test_registry_dispatches_by_mimetype(self):
        from beever_atlas.services.media_extractors import (
            AudioExtractor,
            VideoExtractor,
            create_default_registry,
        )

        registry = create_default_registry()
        assert isinstance(
            registry.get_extractor(mimetype="video/mp4"), VideoExtractor
        )
        assert isinstance(
            registry.get_extractor(mimetype="audio/mpeg"), AudioExtractor
        )

    def test_registry_unknown_type_returns_none(self):
        from beever_atlas.services.media_extractors import create_default_registry

        registry = create_default_registry()
        assert registry.get_extractor(filename="data.xyz", mimetype="application/octet-stream") is None

    @pytest.mark.asyncio
    async def test_registry_fallback_for_unknown(self):
        from beever_atlas.services.media_extractors import create_default_registry

        registry = create_default_registry()
        result = await registry.extract(b"data", "file.xyz")
        assert "file.xyz" in result.text
        assert result.metadata.get("fallback") is True

    @pytest.mark.asyncio
    async def test_pdf_extractor(self):
        from beever_atlas.services.media_extractors import PdfExtractor

        extractor = PdfExtractor()
        # Create a minimal valid PDF
        try:
            from pypdf import PdfWriter
            writer = PdfWriter()
            writer.add_blank_page(width=72, height=72)
            buf = io.BytesIO()
            writer.write(buf)
            pdf_bytes = buf.getvalue()
        except ImportError:
            pytest.skip("pypdf not installed")

        result = await extractor.extract(pdf_bytes, "test.pdf")
        assert result.media_type == "pdf"
        assert "test.pdf" in result.text

    @pytest.mark.asyncio
    async def test_office_extractor_docx(self):
        """Test docx extraction (requires python-docx)."""
        try:
            from docx import Document
        except ImportError:
            pytest.skip("python-docx not installed")

        from beever_atlas.services.media_extractors import OfficeExtractor

        # Create a minimal docx
        doc = Document()
        doc.add_heading("Test Heading", level=1)
        doc.add_paragraph("Test paragraph content.")
        buf = io.BytesIO()
        doc.save(buf)
        docx_bytes = buf.getvalue()

        extractor = OfficeExtractor()
        result = await extractor.extract(docx_bytes, "test.docx")
        assert "Test Heading" in result.text
        assert "Test paragraph" in result.text

    @pytest.mark.asyncio
    async def test_office_extractor_xlsx(self):
        """Test xlsx extraction (requires openpyxl)."""
        try:
            from openpyxl import Workbook
        except ImportError:
            pytest.skip("openpyxl not installed")

        from beever_atlas.services.media_extractors import OfficeExtractor

        wb = Workbook()
        ws = wb.active
        ws.title = "Data"
        ws["A1"] = "Name"
        ws["B1"] = "Value"
        ws["A2"] = "Test"
        ws["B2"] = 42
        buf = io.BytesIO()
        wb.save(buf)
        xlsx_bytes = buf.getvalue()

        extractor = OfficeExtractor()
        result = await extractor.extract(xlsx_bytes, "data.xlsx")
        assert "Sheet: Data" in result.text
        assert "Name" in result.text

    @pytest.mark.asyncio
    async def test_office_extractor_pptx(self):
        """Test pptx extraction (requires python-pptx)."""
        try:
            from pptx import Presentation
        except ImportError:
            pytest.skip("python-pptx not installed")

        from beever_atlas.services.media_extractors import OfficeExtractor

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Slide Title"
        buf = io.BytesIO()
        prs.save(buf)
        pptx_bytes = buf.getvalue()

        extractor = OfficeExtractor()
        result = await extractor.extract(pptx_bytes, "slides.pptx")
        assert "Slide 1" in result.text
        assert "Slide Title" in result.text


# ── Group 4: Semantic Search ──────────────────────────────────────────


class TestSemanticSearch:
    """Tests for Weaviate semantic search methods."""

    def test_weaviate_store_has_semantic_search(self):
        from beever_atlas.stores.weaviate_store import WeaviateStore

        store = WeaviateStore("http://localhost:8080")
        assert hasattr(store, "semantic_search")
        assert hasattr(store, "hybrid_search")
        assert hasattr(store, "supersede_fact")

    def test_supersession_properties_in_schema(self):
        from beever_atlas.stores.weaviate_store import WeaviateStore

        prop_names = [name for name, _ in WeaviateStore._EXPECTED_PROPERTIES]
        assert "superseded_by" in prop_names
        assert "supersedes" in prop_names
        assert "potential_contradiction" in prop_names


# ── Group 5: Temporal Fact Lifecycle ──────────────────────────────────


class TestTemporalFactLifecycle:
    """Tests for contradiction detection and fact supersession."""

    def test_atomic_fact_has_supersession_fields(self):
        from beever_atlas.models import AtomicFact

        fact = AtomicFact(memory_text="test")
        assert fact.superseded_by is None
        assert fact.supersedes is None
        assert fact.potential_contradiction is False

    def test_atomic_fact_supersession_roundtrip(self):
        from beever_atlas.models import AtomicFact

        fact = AtomicFact(
            memory_text="test",
            superseded_by="uuid-123",
            supersedes="uuid-456",
            potential_contradiction=True,
        )
        assert fact.superseded_by == "uuid-123"
        assert fact.supersedes == "uuid-456"
        assert fact.potential_contradiction is True

    @pytest.mark.asyncio
    async def test_detect_contradictions_empty(self):
        from beever_atlas.models import AtomicFact
        from beever_atlas.services.contradiction_detector import detect_contradictions

        new_fact = AtomicFact(memory_text="test")
        result = await detect_contradictions(new_fact, [])
        assert result == []


# ── Group 6: Cross-Batch Thread Context ───────────────────────────────


class TestCrossBatchThreadContext:
    """Tests for cross-batch thread context resolution."""

    def test_config_has_thread_context_settings(self):
        from beever_atlas.infra.config import Settings

        s = Settings(google_api_key="x")
        assert s.cross_batch_thread_context_enabled is True
        assert s.thread_context_max_length == 200


# ── Group 7: Soft Orphan Handling ─────────────────────────────────────


class TestSoftOrphanHandling:
    """Tests for pending entity status and promotion."""

    def test_graph_entity_has_status_fields(self):
        from beever_atlas.models import GraphEntity

        entity = GraphEntity(name="Test", type="Project")
        assert entity.status == "active"
        assert entity.pending_since is None

    def test_graph_entity_pending_status(self):
        from beever_atlas.models import GraphEntity

        entity = GraphEntity(
            name="Test",
            type="Project",
            status="pending",
            pending_since=datetime.now(tz=UTC),
        )
        assert entity.status == "pending"
        assert entity.pending_since is not None

    def test_graph_entity_name_vector(self):
        from beever_atlas.models import GraphEntity

        entity = GraphEntity(name="Test", type="Project", name_vector=[0.1, 0.2])
        assert entity.name_vector == [0.1, 0.2]

    def test_neo4j_store_has_orphan_methods(self):
        from beever_atlas.stores.neo4j_store import Neo4jStore

        store = Neo4jStore("bolt://localhost:7687", "neo4j", "test")
        assert hasattr(store, "promote_pending_entity")
        assert hasattr(store, "prune_expired_pending")

    def test_config_has_orphan_settings(self):
        from beever_atlas.infra.config import Settings

        s = Settings(google_api_key="x")
        assert s.orphan_grace_period_days == 7


# ── Integration: Config Completeness ──────────────────────────────────


class TestConfigCompleteness:
    """Verify all new configuration settings exist."""

    def test_all_new_settings_exist(self):
        from beever_atlas.infra.config import Settings

        s = Settings(google_api_key="x")
        # Coreference
        assert s.coref_enabled is True
        assert s.coref_history_limit == 20
        assert s.coref_model == "gemini-2.5-flash"
        # Semantic dedup
        assert s.entity_similarity_threshold == 0.85
        assert s.merge_rejection_ttl_days == 30
        # Multimodal
        assert s.media_video_max_duration_minutes == 10
        assert s.media_video_max_size_mb == 100
        assert s.media_audio_max_duration_minutes == 30
        assert s.media_office_max_chars == 10000
        # Semantic search
        assert s.semantic_search_min_similarity == 0.7
        # Temporal lifecycle
        assert s.contradiction_confidence_threshold == 0.8
        assert s.contradiction_flag_threshold == 0.5
        # Thread context
        assert s.cross_batch_thread_context_enabled is True
        assert s.thread_context_max_length == 200
        # Orphan
        assert s.orphan_grace_period_days == 7
