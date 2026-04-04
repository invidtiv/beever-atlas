"""Dry-run script for testing pipeline hardening features.

Tests all 7 hardening features in isolation (no live services required):
  1. Coreference resolution (pronoun detection + LLM skip)
  2. Semantic entity deduplication (embedding similarity)
  3. Multimodal expansion (registry dispatch + office extraction)
  4. Semantic search (method existence + schema)
  5. Temporal fact lifecycle (model fields + contradiction detection)
  6. Cross-batch thread context (config + resolution logic)
  7. Soft orphan handling (model fields + status management)

Usage:
    python scripts/dry_run_pipeline_hardening.py
"""

from __future__ import annotations

import asyncio
import io
import sys
from datetime import UTC, datetime
from unittest.mock import AsyncMock, MagicMock


def _header(title: str) -> None:
    print(f"\n{'=' * 60}")
    print(f"  {title}")
    print(f"{'=' * 60}")


def _ok(msg: str) -> None:
    print(f"  [PASS] {msg}")


def _fail(msg: str) -> None:
    print(f"  [FAIL] {msg}")


def _skip(msg: str) -> None:
    print(f"  [SKIP] {msg}")


passed = 0
failed = 0
skipped = 0


def check(condition: bool, msg: str) -> None:
    global passed, failed
    if condition:
        _ok(msg)
        passed += 1
    else:
        _fail(msg)
        failed += 1


def skip(msg: str) -> None:
    global skipped
    _skip(msg)
    skipped += 1


async def test_group1_coreference() -> None:
    _header("Group 1: Coreference Resolution")

    from beever_atlas.services.coreference_resolver import (
        has_resolvable_references,
        resolve_coreferences,
    )

    # Test pronoun detection
    check(
        has_resolvable_references([{"text": "It uses Redis"}]),
        "Detects pronoun 'it'",
    )
    check(
        has_resolvable_references([{"text": "They approved the plan"}]),
        "Detects pronoun 'they'",
    )
    check(
        has_resolvable_references([{"text": "The project needs work"}]),
        "Detects implicit reference 'the project'",
    )
    check(
        not has_resolvable_references([{"text": "PostgreSQL 15 released"}]),
        "No false positive on plain text",
    )

    # Test no-pronoun skip
    result = await resolve_coreferences([{"text": "PostgreSQL 15 released"}])
    check(
        result[0].get("raw_text") == "PostgreSQL 15 released",
        "Preserves raw_text when no pronouns",
    )

    # Test empty batch
    result = await resolve_coreferences([])
    check(result == [], "Empty batch returns empty")


async def test_group2_semantic_dedup() -> None:
    _header("Group 2: Semantic Entity Deduplication")

    from beever_atlas.stores.entity_registry import EntityRegistry

    # Test embedding similarity
    mock_neo4j = AsyncMock()
    mock_neo4j.execute_query = AsyncMock(return_value=[
        {"name": "Beever Atlas", "vec": [0.9, 0.1, 0.0]},
        {"name": "Redis", "vec": [0.1, 0.9, 0.0]},
    ])
    registry = EntityRegistry(mock_neo4j)

    results = await registry.find_similar_by_embedding(
        name="Atlas", name_vector=[0.85, 0.15, 0.0], threshold=0.8
    )
    similar_names = [r[0] for r in results]
    check("Beever Atlas" in similar_names, "Finds semantically similar 'Beever Atlas'")
    check("Redis" not in similar_names, "Excludes dissimilar 'Redis'")

    # Test rejection cache
    check(not registry.is_merge_rejected("A", "B"), "No rejection initially")
    registry.cache_merge_rejection("A", "B")
    check(registry.is_merge_rejected("A", "B"), "Rejection cached")
    check(registry.is_merge_rejected("B", "A"), "Rejection is order-independent")


async def test_group3_multimodal() -> None:
    _header("Group 3: Multimodal Expansion")

    from beever_atlas.services.media_extractors import (
        AudioExtractor,
        ImageExtractor,
        OfficeExtractor,
        PdfExtractor,
        VideoExtractor,
        create_default_registry,
    )

    registry = create_default_registry()

    # Test registry dispatch
    check(isinstance(registry.get_extractor(filename="a.pdf"), PdfExtractor), "PDF dispatch")
    check(isinstance(registry.get_extractor(filename="a.png"), ImageExtractor), "Image dispatch")
    check(isinstance(registry.get_extractor(filename="a.docx"), OfficeExtractor), "DOCX dispatch")
    check(isinstance(registry.get_extractor(filename="a.xlsx"), OfficeExtractor), "XLSX dispatch")
    check(isinstance(registry.get_extractor(filename="a.pptx"), OfficeExtractor), "PPTX dispatch")
    check(isinstance(registry.get_extractor(mimetype="video/mp4"), VideoExtractor), "Video dispatch")
    check(isinstance(registry.get_extractor(mimetype="audio/mpeg"), AudioExtractor), "Audio dispatch")
    check(registry.get_extractor(filename="a.xyz") is None, "Unknown returns None")

    # Test fallback
    result = await registry.extract(b"data", "file.xyz")
    check("file.xyz" in result.text, "Fallback includes filename")
    check(result.metadata.get("fallback") is True, "Fallback flagged in metadata")

    # Test docx extraction
    try:
        from docx import Document
        doc = Document()
        doc.add_heading("Quarterly Report", level=1)
        doc.add_paragraph("Revenue increased by 15%.")
        buf = io.BytesIO()
        doc.save(buf)

        extractor = OfficeExtractor()
        result = await extractor.extract(buf.getvalue(), "report.docx")
        check("Quarterly Report" in result.text, "DOCX heading extracted")
        check("Revenue" in result.text, "DOCX paragraph extracted")
    except ImportError:
        skip("python-docx not installed — skipping docx extraction test")

    # Test xlsx extraction
    try:
        from openpyxl import Workbook
        wb = Workbook()
        ws = wb.active
        if ws is not None:
            ws.title = "Sales"
            ws["A1"] = "Product"
            ws["B1"] = "Revenue"
        buf = io.BytesIO()
        wb.save(buf)

        extractor = OfficeExtractor()
        result = await extractor.extract(buf.getvalue(), "sales.xlsx")
        check("Sheet: Sales" in result.text, "XLSX sheet name extracted")
        check("Product" in result.text, "XLSX cell data extracted")
    except ImportError:
        skip("openpyxl not installed — skipping xlsx extraction test")

    # Test pptx extraction
    try:
        from pptx import Presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Q4 Review"
        buf = io.BytesIO()
        prs.save(buf)

        extractor = OfficeExtractor()
        result = await extractor.extract(buf.getvalue(), "deck.pptx")
        check("Q4 Review" in result.text, "PPTX slide text extracted")
    except ImportError:
        skip("python-pptx not installed — skipping pptx extraction test")


async def test_group4_semantic_search() -> None:
    _header("Group 4: Semantic Search Activation")

    from beever_atlas.stores.weaviate_store import WeaviateStore

    store = WeaviateStore("http://localhost:8080")
    check(hasattr(store, "semantic_search"), "WeaviateStore has semantic_search")
    check(hasattr(store, "hybrid_search"), "WeaviateStore has hybrid_search")
    check(hasattr(store, "supersede_fact"), "WeaviateStore has supersede_fact")
    check(hasattr(store, "flag_potential_contradiction"), "WeaviateStore has flag_potential_contradiction")

    # Check schema includes new properties
    prop_names = [name for name, _ in WeaviateStore._EXPECTED_PROPERTIES]
    check("superseded_by" in prop_names, "Schema has superseded_by")
    check("supersedes" in prop_names, "Schema has supersedes")
    check("potential_contradiction" in prop_names, "Schema has potential_contradiction")


async def test_group5_temporal_lifecycle() -> None:
    _header("Group 5: Temporal Fact Lifecycle")

    from beever_atlas.models import AtomicFact
    from beever_atlas.services.contradiction_detector import detect_contradictions

    # Test model fields
    fact = AtomicFact(memory_text="test")
    check(fact.superseded_by is None, "Default superseded_by is None")
    check(fact.supersedes is None, "Default supersedes is None")
    check(fact.potential_contradiction is False, "Default potential_contradiction is False")

    # Test supersession roundtrip
    fact2 = AtomicFact(
        memory_text="test",
        superseded_by="id-1",
        supersedes="id-2",
        potential_contradiction=True,
    )
    check(fact2.superseded_by == "id-1", "Superseded_by field roundtrips")
    check(fact2.potential_contradiction is True, "Potential_contradiction field roundtrips")

    # Test contradiction detection with empty candidates
    result = await detect_contradictions(fact, [])
    check(result == [], "No contradictions when no candidates")


async def test_group6_thread_context() -> None:
    _header("Group 6: Cross-Batch Thread Context")

    from beever_atlas.infra.config import Settings

    s = Settings(google_api_key="test")
    check(s.cross_batch_thread_context_enabled is True, "Thread context enabled by default")
    check(s.thread_context_max_length == 200, "Max length is 200")


async def test_group7_soft_orphan() -> None:
    _header("Group 7: Soft Orphan Handling")

    from beever_atlas.models import GraphEntity
    from beever_atlas.stores.neo4j_store import Neo4jStore

    # Test model fields
    entity = GraphEntity(name="Project X", type="Project")
    check(entity.status == "active", "Default status is active")
    check(entity.pending_since is None, "Default pending_since is None")
    check(entity.name_vector is None, "Default name_vector is None")

    # Test pending entity
    pending = GraphEntity(
        name="New Thing",
        type="Project",
        status="pending",
        pending_since=datetime.now(tz=UTC),
    )
    check(pending.status == "pending", "Pending status set correctly")
    check(pending.pending_since is not None, "Pending_since timestamp set")

    # Test store has methods
    store = Neo4jStore("bolt://localhost:7687", "neo4j", "test")
    check(hasattr(store, "promote_pending_entity"), "Neo4j has promote_pending_entity")
    check(hasattr(store, "prune_expired_pending"), "Neo4j has prune_expired_pending")

    # Test config
    from beever_atlas.infra.config import Settings
    s = Settings(google_api_key="test")
    check(s.orphan_grace_period_days == 7, "Grace period default is 7 days")


async def test_phase1_backend_wiring() -> None:
    _header("Phase 1: Backend Pipeline Wiring (v2)")

    # 1.1: ExtractedEntity has status field
    from beever_atlas.agents.schemas.extraction import ExtractedEntity
    e = ExtractedEntity(name="Test", type="Project")
    check(e.status == "active", "ExtractedEntity default status is 'active'")
    e2 = ExtractedEntity(name="Orphan", type="Project", status="pending")
    check(e2.status == "pending", "ExtractedEntity accepts 'pending' status")

    # 1.1: Validator prompt mentions 'pending'
    from beever_atlas.agents.prompts.cross_batch_validator import CROSS_BATCH_VALIDATOR_INSTRUCTION
    check("pending" in CROSS_BATCH_VALIDATOR_INSTRUCTION, "Validator prompt mentions 'pending'")
    check("Do NOT remove" in CROSS_BATCH_VALIDATOR_INSTRUCTION, "Validator prompt says do NOT remove orphans")
    check('"status"' in CROSS_BATCH_VALIDATOR_INSTRUCTION, "Validator output format includes status field")

    # 1.2: EntityRegistry has batch embedding method
    from beever_atlas.stores.entity_registry import EntityRegistry
    registry = EntityRegistry(MagicMock())
    check(hasattr(registry, "compute_name_embeddings_batch"), "EntityRegistry has batch embedding method")

    # 1.3: Contradiction detector output schema exists
    from beever_atlas.agents.schemas.validation import ContradictionResult, ContradictionReport
    cr = ContradictionResult(existing_fact_id="abc", confidence=0.9, reason="test")
    check(cr.confidence == 0.9, "ContradictionResult model works")
    report = ContradictionReport()
    check(report.contradictions == [], "ContradictionReport defaults to empty")


async def test_phase3_gemini_flash() -> None:
    _header("Phase 3: Gemini Flash Media (v2)")

    from beever_atlas.services.media_extractors import VideoExtractor, AudioExtractor
    import inspect

    # Check that Whisper references are gone from these extractors
    video_src = inspect.getsource(VideoExtractor._transcribe_audio)
    check("whisper" not in video_src.lower(), "VideoExtractor no longer references Whisper")
    check("genai" in video_src or "google" in video_src, "VideoExtractor uses Google/genai")

    audio_src = inspect.getsource(AudioExtractor.extract)
    check("whisper" not in audio_src.lower(), "AudioExtractor no longer references Whisper")
    check("genai" in audio_src or "google" in audio_src, "AudioExtractor uses Google/genai")


async def test_phase4_docs() -> None:
    _header("Phase 4: Documentation & Schema (v2)")

    import os
    doc_path = os.path.join(os.path.dirname(__file__), "..", "docs", "pipeline-architecture.md")
    check(os.path.exists(doc_path), "Pipeline architecture doc exists")

    with open(doc_path) as f:
        content = f.read()
    check("Stage 1: Preprocessor" in content, "Doc covers Stage 1")
    check("Coreference Resolution" in content, "Doc covers coreference")
    check("Entity Lifecycle" in content, "Doc covers entity lifecycle")
    check("Gemini Flash" in content, "Doc mentions Gemini Flash")
    check("Contradiction Detection" in content, "Doc covers contradiction detection")


async def main() -> None:
    print("\n" + "=" * 60)
    print("  PIPELINE HARDENING DRY RUN (v2)")
    print("  Testing all 7 feature groups + Phase 1-4 integrations")
    print("=" * 60)

    await test_group1_coreference()
    await test_group2_semantic_dedup()
    await test_group3_multimodal()
    await test_group4_semantic_search()
    await test_group5_temporal_lifecycle()
    await test_group6_thread_context()
    await test_group7_soft_orphan()
    await test_phase1_backend_wiring()
    await test_phase3_gemini_flash()
    await test_phase4_docs()

    print(f"\n{'=' * 60}")
    print(f"  RESULTS: {passed} passed, {failed} failed, {skipped} skipped")
    print(f"{'=' * 60}\n")

    if failed > 0:
        sys.exit(1)


if __name__ == "__main__":
    asyncio.run(main())
