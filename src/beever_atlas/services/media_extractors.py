"""Media extractor registry and extractors for multimodal content.

Provides a registry-based dispatch system for media file processing,
replacing hardcoded if/else branches. Each extractor handles specific
MIME types and returns extracted text content.
"""

from __future__ import annotations

import asyncio
import io
import logging
import re
from abc import ABC, abstractmethod
from dataclasses import dataclass, field
from typing import Any

from beever_atlas.agents.prompts.media import (
    AUDIO_TRANSCRIPTION_PROMPT,
    DOCUMENT_DIGEST_PROMPT,
    IMAGE_DESCRIPTION_PROMPT,
    VIDEO_ANALYSIS_PROMPT,
)
from beever_atlas.infra.config import get_settings
from beever_atlas.infra.rate_limit import GEMINI_LIMITER

logger = logging.getLogger(__name__)

# ── Shared Gemini client ──────────────────────────────────────────────
_gemini_client: Any = None
_gemini_client_lock: asyncio.Lock | None = None

# Bounds concurrent in-flight Gemini image calls across all coroutines.
# Set to 4 so up to 4 images are described in parallel; the outer
# GEMINI_LIMITER still enforces the RPM budget.
_IMAGE_SEMAPHORE: asyncio.Semaphore | None = None
_IMAGE_SEMAPHORE_LOCK: asyncio.Lock | None = None


def _get_image_semaphore() -> asyncio.Semaphore:
    global _IMAGE_SEMAPHORE, _IMAGE_SEMAPHORE_LOCK
    # Fast path — already created
    if _IMAGE_SEMAPHORE is not None:
        return _IMAGE_SEMAPHORE
    # Lazy-init the lock itself (must happen inside a running loop)
    if _IMAGE_SEMAPHORE_LOCK is None:
        _IMAGE_SEMAPHORE_LOCK = asyncio.Lock()
    # Caller must be inside an async context; return existing or create
    concurrency = get_settings().image_extractor_concurrency
    _IMAGE_SEMAPHORE = asyncio.Semaphore(concurrency)
    return _IMAGE_SEMAPHORE


def _get_client_lock() -> asyncio.Lock:
    global _gemini_client_lock
    if _gemini_client_lock is None:
        _gemini_client_lock = asyncio.Lock()
    return _gemini_client_lock


async def _get_gemini_client() -> Any:
    """Return a shared genai.Client, creating it lazily under a lock."""
    global _gemini_client
    lock = _get_client_lock()
    async with lock:
        if _gemini_client is None:
            from google import genai

            settings = get_settings()
            _gemini_client = genai.Client(api_key=settings.google_api_key)
        return _gemini_client


async def _poll_file_active(
    client: Any,
    file_name: str,
    *,
    poll_interval: float = 2.0,
    max_wait: float = 90.0,
) -> None:
    """Poll Gemini Files API until the uploaded file reaches ACTIVE state."""
    elapsed = 0.0
    while elapsed < max_wait:
        file_info = await client.aio.files.get(name=file_name)
        state = file_info.state
        if state == "ACTIVE":
            return
        if state == "FAILED":
            error_msg = getattr(file_info, "error", None) or "unknown error"
            raise RuntimeError(f"File processing failed: {error_msg}")
        await asyncio.sleep(poll_interval)
        elapsed += poll_interval

    raise TimeoutError(f"File {file_name} did not reach ACTIVE state within {max_wait}s")


@dataclass
class MediaContent:
    """Result of media extraction."""

    text: str = ""
    media_urls: list[str] = field(default_factory=list)
    media_type: str = ""
    metadata: dict[str, Any] = field(default_factory=dict)
    chunks: list[str] = field(default_factory=list)


class MediaExtractor(ABC):
    """Base class for media extractors."""

    @property
    @abstractmethod
    def supported_mime_types(self) -> list[str]:
        """MIME types this extractor handles."""

    @property
    @abstractmethod
    def supported_extensions(self) -> list[str]:
        """File extensions this extractor handles (without dot)."""

    @abstractmethod
    async def extract(
        self, data: bytes, filename: str, metadata: dict[str, Any] | None = None
    ) -> MediaContent:
        """Extract content from file bytes."""

    async def _digest_document(self, text: str) -> str:
        """Digest a document into a concise summary using direct Gemini API."""
        settings = get_settings()
        if not settings.google_api_key:
            return text[:4000] + ("\n[...truncated]" if len(text) > 4000 else "")
        try:
            from google.genai import types as genai_types

            client = await _get_gemini_client()

            prompt = DOCUMENT_DIGEST_PROMPT.format(document_text=text[:8000])

            async with GEMINI_LIMITER:
                response = await asyncio.wait_for(
                    client.aio.models.generate_content(
                        model=settings.media_vision_model,
                        contents=[
                            genai_types.Content(
                                role="user",
                                parts=[genai_types.Part.from_text(text=prompt)],
                            )
                        ],
                    ),
                    timeout=60,
                )
            return response.text or text[:4000]
        except Exception:
            logger.warning("Document digestion failed", exc_info=True)
            return text[:4000] + ("\n[...truncated]" if len(text) > 4000 else "")


class MediaExtractorRegistry:
    """Registry that dispatches file processing to the appropriate extractor."""

    def __init__(self) -> None:
        self._mime_map: dict[str, MediaExtractor] = {}
        self._ext_map: dict[str, MediaExtractor] = {}

    def register(self, extractor: MediaExtractor) -> None:
        """Register an extractor for its supported MIME types and extensions."""
        for mime in extractor.supported_mime_types:
            self._mime_map[mime.lower()] = extractor
        for ext in extractor.supported_extensions:
            self._ext_map[ext.lower()] = extractor

    def get_extractor(self, mimetype: str = "", filename: str = "") -> MediaExtractor | None:
        """Find an extractor by MIME type or file extension."""
        if mimetype:
            extractor = self._mime_map.get(mimetype.lower())
            if extractor:
                return extractor
        if filename and "." in filename:
            ext = filename.rsplit(".", 1)[-1].lower()
            return self._ext_map.get(ext)
        return None

    async def extract(
        self,
        data: bytes,
        filename: str,
        mimetype: str = "",
        metadata: dict[str, Any] | None = None,
    ) -> MediaContent:
        """Dispatch to the appropriate extractor, or return metadata-only fallback."""
        extractor = self.get_extractor(mimetype, filename)
        if extractor is None:
            ext = filename.rsplit(".", 1)[-1] if "." in filename else "unknown"
            return MediaContent(
                text=f"[Attachment: {filename} ({ext})]",
                media_type=ext,
                metadata={"fallback": True},
            )
        return await extractor.extract(data, filename, metadata)


# ── Concrete Extractors ────────────────────────────────────────────────


class PdfExtractor(MediaExtractor):
    """Extract text from PDF files using pypdf with page-chunked output."""

    @property
    def supported_mime_types(self) -> list[str]:
        return ["application/pdf"]

    @property
    def supported_extensions(self) -> list[str]:
        return ["pdf"]

    async def extract(
        self, data: bytes, filename: str, metadata: dict[str, Any] | None = None
    ) -> MediaContent:
        from beever_atlas.infra.config import get_settings

        settings = get_settings()

        pages = await asyncio.to_thread(self._extract_pages, data, settings.pdf_max_pages)
        size_kb = len(data) // 1024

        if not any(p.strip() for p in pages):
            return MediaContent(
                text=f"[Attachment: {filename} (PDF, {size_kb} kB, no extractable text)]",
                media_type="pdf",
            )

        full_text = "\n\n".join(p.strip() for p in pages if p.strip())
        header = f"[Attachment: {filename} (PDF, {size_kb} kB, {len(pages)} pages)]"

        if settings.media_digest_enabled:
            digest_content = await self._digest_document(full_text)
            desc = f"{header}\n[Document Digest]:\n{digest_content}"
        else:
            # Skip LLM digestion — return truncated raw text for speed
            truncated = full_text[:4000]
            if len(full_text) > 4000:
                truncated += "\n[...truncated]"
            desc = f"{header}\n{truncated}"

        return MediaContent(
            text=desc,
            media_type="pdf",
        )

    @staticmethod
    def _extract_pages(data: bytes, max_pages: int) -> list[str]:
        """Extract text from each page independently. Returns one string per page."""
        try:
            from pypdf import PdfReader

            reader = PdfReader(io.BytesIO(data))
            pages: list[str] = []
            total = len(reader.pages)
            limit = min(total, max_pages)

            for i in range(limit):
                page_text = reader.pages[i].extract_text() or ""
                pages.append(page_text.strip())

            if total > max_pages:
                remaining = total - max_pages
                pages.append(f"[...remaining {remaining} pages not processed]")

            return pages
        except Exception:
            logger.warning("PdfExtractor: page extraction failed", exc_info=True)
            return []


class ImageExtractor(MediaExtractor):
    """Describe images using Gemini vision API."""

    # Patterns suggesting visual content worth describing
    _VISUAL_RE = re.compile(
        r"screenshot|diagram|chart|graph|whiteboard|mockup|wireframe|design|sketch",
        re.IGNORECASE,
    )
    _ATTACHMENT_REF_RE = re.compile(
        r"see attached|check this|look at this|attached|screenshot|see above|here'?s the",
        re.IGNORECASE,
    )

    @property
    def supported_mime_types(self) -> list[str]:
        return [
            "image/png",
            "image/jpeg",
            "image/jpg",
            "image/gif",
            "image/webp",
        ]

    @property
    def supported_extensions(self) -> list[str]:
        return ["png", "jpg", "jpeg", "gif", "webp"]

    async def extract(
        self, data: bytes, filename: str, metadata: dict[str, Any] | None = None
    ) -> MediaContent:
        message_text = (metadata or {}).get("message_text", "")
        size_kb = len(data) // 1024

        if not self._should_use_vision(message_text, filename):
            return MediaContent(
                text=f"[Attachment: {filename} (image)]",
                media_type="image",
            )

        description = await self._describe_image(data, message_text, filename)
        if description:
            desc = (
                f"[Attachment: {filename} (image, {size_kb} kB)]\n"
                f"[Image description]: {description}"
            )
        else:
            desc = f"[Attachment: {filename} (image, {size_kb} kB)]"
        return MediaContent(text=desc, media_type="image")

    def _should_use_vision(self, message_text: str, filename: str) -> bool:
        text = (message_text or "").strip()
        if len(text) < 50:
            return True
        if self._ATTACHMENT_REF_RE.search(text):
            return True
        if filename and self._VISUAL_RE.search(filename):
            return True
        return False

    _IMAGE_MIME_MAP: dict[str, str] = {
        "png": "image/png",
        "jpg": "image/jpeg",
        "jpeg": "image/jpeg",
        "gif": "image/gif",
        "webp": "image/webp",
    }

    async def _describe_image(self, data: bytes, message_context: str, filename: str = "") -> str:
        """Describe an image using direct Gemini API with multimodal parts."""
        settings = get_settings()
        if not settings.google_api_key:
            return ""
        # Infer mime type from filename extension
        ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else "png"
        mime_type = self._IMAGE_MIME_MAP.get(ext, "image/png")

        try:
            from google.genai import types as genai_types

            client = await _get_gemini_client()

            prompt = IMAGE_DESCRIPTION_PROMPT
            if message_context:
                prompt += f"\n\nMessage context: {message_context[:200]}"

            logger.info(
                "ImageExtractor: calling generate_content for %s (%s, %d bytes)",
                filename,
                mime_type,
                len(data),
            )
            async with _get_image_semaphore():
                async with GEMINI_LIMITER:
                    response = await asyncio.wait_for(
                        client.aio.models.generate_content(
                            model=settings.media_vision_model,
                            contents=[
                                genai_types.Content(
                                    role="user",
                                    parts=[
                                        genai_types.Part.from_bytes(
                                            data=data,
                                            mime_type=mime_type,
                                        ),
                                        genai_types.Part.from_text(text=prompt),
                                    ],
                                )
                            ],
                        ),
                        timeout=60,
                    )
            result = response.text or ""
            logger.info(
                "ImageExtractor: result for %s — %d chars",
                filename,
                len(result),
            )
            return result
        except Exception:
            logger.error(
                "ImageExtractor: vision description failed for %s (mime=%s)",
                filename,
                mime_type,
                exc_info=True,
            )
            return ""


class OfficeExtractor(MediaExtractor):
    """Extract text from Office documents (docx, xlsx, pptx)."""

    @property
    def supported_mime_types(self) -> list[str]:
        return [
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "application/msword",
            "application/vnd.ms-excel",
            "application/vnd.ms-powerpoint",
        ]

    @property
    def supported_extensions(self) -> list[str]:
        return ["docx", "xlsx", "pptx", "doc", "xls", "ppt"]

    async def extract(
        self, data: bytes, filename: str, metadata: dict[str, Any] | None = None
    ) -> MediaContent:
        ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
        settings = get_settings()
        max_chars = settings.media_office_max_chars
        size_kb = len(data) // 1024

        text = await asyncio.to_thread(self._extract_text, data, ext, max_chars)

        if text.strip():
            if settings.media_digest_enabled:
                digest_content = await self._digest_document(text)
                desc = (
                    f"[Attachment: {filename} ({ext.upper()}, {size_kb} kB)]\n"
                    f"[Document Digest]:\n{digest_content}"
                )
            else:
                truncated = text[:4000]
                if len(text) > 4000:
                    truncated += "\n[...truncated]"
                desc = f"[Attachment: {filename} ({ext.upper()}, {size_kb} kB)]\n{truncated}"
        else:
            desc = f"[Attachment: {filename} ({ext.upper()}, {size_kb} kB, no extractable text)]"

        return MediaContent(text=desc, media_type=ext)

    @staticmethod
    def _extract_text(data: bytes, ext: str, max_chars: int) -> str:
        try:
            if ext in ("docx", "doc"):
                return OfficeExtractor._extract_docx(data, max_chars)
            elif ext in ("xlsx", "xls"):
                return OfficeExtractor._extract_xlsx(data, max_chars)
            elif ext in ("pptx", "ppt"):
                return OfficeExtractor._extract_pptx(data, max_chars)
        except Exception:
            logger.warning("OfficeExtractor: extraction failed for .%s", ext, exc_info=True)
        return ""

    @staticmethod
    def _extract_docx(data: bytes, max_chars: int) -> str:
        from docx import Document

        doc = Document(io.BytesIO(data))
        parts: list[str] = []
        char_count = 0
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            if para.style and para.style.name and para.style.name.startswith("Heading"):
                parts.append(f"## {text}")
            else:
                parts.append(text)
            char_count += len(text)
            if char_count >= max_chars:
                parts.append("[...truncated]")
                break
        return "\n".join(parts)

    @staticmethod
    def _extract_xlsx(data: bytes, max_chars: int) -> str:
        from openpyxl import load_workbook

        wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        parts: list[str] = []
        char_count = 0
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            parts.append(f"Sheet: {sheet_name}")
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(cells):
                    line = " | ".join(cells)
                    parts.append(line)
                    char_count += len(line)
                    if char_count >= max_chars:
                        parts.append("[...truncated]")
                        break
            if char_count >= max_chars:
                break
        wb.close()
        return "\n".join(parts)

    @staticmethod
    def _extract_pptx(data: bytes, max_chars: int) -> str:
        from pptx import Presentation

        prs = Presentation(io.BytesIO(data))
        parts: list[str] = []
        char_count = 0
        for i, slide in enumerate(prs.slides, 1):
            slide_text: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_text.append(text)
            notes = ""
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()

            content = f"Slide {i}: " + " ".join(slide_text)
            if notes:
                content += f"\nNotes: {notes}"
            parts.append(content)
            char_count += len(content)
            if char_count >= max_chars:
                parts.append("[...truncated]")
                break
        return "\n".join(parts)


class VideoExtractor(MediaExtractor):
    """Extract content from video files via combined audio + visual analysis."""

    @property
    def supported_mime_types(self) -> list[str]:
        return ["video/mp4", "video/quicktime", "video/x-msvideo", "video/webm", "video/mpeg"]

    @property
    def supported_extensions(self) -> list[str]:
        return ["mp4", "mov", "avi", "webm", "mpeg", "mkv"]

    async def extract(
        self, data: bytes, filename: str, metadata: dict[str, Any] | None = None
    ) -> MediaContent:
        settings = get_settings()
        size_mb = len(data) / (1024 * 1024)
        max_size = settings.media_video_max_size_mb

        if size_mb > max_size:
            return MediaContent(
                text=f"[Attachment: {filename} (video, {size_mb:.1f} MB — exceeds {max_size} MB limit)]",
                media_type="video",
            )

        parts: list[str] = [f"[Attachment: {filename} (video, {size_mb:.1f} MB)]"]

        # Single Gemini call for combined transcript + visual analysis
        analysis = await self._analyze_video(data, filename)
        if analysis:
            parts.append(f"[Video summary]: {analysis}")

        return MediaContent(text="\n".join(parts), media_type="video")

    async def _analyze_video(self, data: bytes, filename: str) -> str:
        """Analyze video using Gemini Files API + content generation."""
        settings = get_settings()
        if not settings.google_api_key:
            logger.debug("VideoExtractor: no Google API key, skipping analysis")
            return ""

        uploaded = None
        client = await _get_gemini_client()
        try:
            from google.genai import types as genai_types

            ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else "mp4"
            mime_map = {
                "mp4": "video/mp4",
                "mov": "video/quicktime",
                "webm": "video/webm",
                "avi": "video/x-msvideo",
            }
            mime_type = mime_map.get(ext, "video/mp4")

            # Upload via Files API
            logger.info(
                "VideoExtractor: uploading %s (%s, %.1f MB)",
                filename,
                mime_type,
                len(data) / (1024 * 1024),
            )
            uploaded = await asyncio.wait_for(
                client.aio.files.upload(
                    file=io.BytesIO(data),
                    config=genai_types.UploadFileConfig(mime_type=mime_type, display_name=filename),
                ),
                timeout=60,
            )
            file_uri: str = uploaded.uri or ""
            logger.info(
                "VideoExtractor: uploaded %s → %s (state=%s)",
                filename,
                uploaded.name,
                uploaded.state,
            )

            # Poll until file is ready for use
            await _poll_file_active(client, uploaded.name)
            logger.info("VideoExtractor: file ACTIVE for %s", filename)

            async with GEMINI_LIMITER:
                response = await asyncio.wait_for(
                    client.aio.models.generate_content(
                        model=settings.media_vision_model,
                        contents=[
                            genai_types.Content(
                                role="user",
                                parts=[
                                    genai_types.Part.from_uri(
                                        file_uri=file_uri,
                                        mime_type=mime_type,
                                    ),
                                    genai_types.Part.from_text(text=VIDEO_ANALYSIS_PROMPT),
                                ],
                            )
                        ],
                    ),
                    timeout=60,
                )
            result = response.text or ""
            logger.info("VideoExtractor: result for %s — %d chars", filename, len(result))
            return result
        except Exception:
            logger.error("VideoExtractor: analysis failed for %s", filename, exc_info=True)
            return ""
        finally:
            if uploaded and getattr(uploaded, "name", None):
                try:
                    await client.aio.files.delete(name=uploaded.name)
                except Exception:
                    logger.debug("VideoExtractor: cleanup failed for %s", filename)

    @staticmethod
    def _parse_analysis(text: str) -> tuple[str, str, str]:
        """Parse structured analysis into (transcript, translation, visual_description)."""
        transcript = ""
        translation = ""
        visual_desc = ""

        current_section = ""
        section_lines: dict[str, list[str]] = {
            "transcript": [],
            "translation": [],
            "visual": [],
        }

        for line in text.split("\n"):
            stripped = line.strip()
            upper = stripped.upper()
            if upper.startswith("TRANSCRIPT:"):
                current_section = "transcript"
                # Capture inline content after the label
                rest = stripped[len("TRANSCRIPT:") :].strip()
                if rest:
                    section_lines["transcript"].append(rest)
            elif upper.startswith("TRANSLATION"):
                current_section = "translation"
                rest = stripped.split(":", 1)[-1].strip() if ":" in stripped else ""
                if rest:
                    section_lines["translation"].append(rest)
            elif upper.startswith("VISUAL DESCRIPTION:") or upper.startswith("VISUAL:"):
                current_section = "visual"
                rest = stripped.split(":", 1)[-1].strip() if ":" in stripped else ""
                if rest:
                    section_lines["visual"].append(rest)
            elif current_section and stripped:
                section_lines[current_section].append(stripped)

        transcript = "\n".join(section_lines["transcript"]).strip()
        translation = "\n".join(section_lines["translation"]).strip()
        visual_desc = "\n".join(section_lines["visual"]).strip()

        # Fallback: if no sections were parsed, treat the whole thing as transcript
        if not transcript and not visual_desc and text.strip():
            transcript = text.strip()

        return transcript, translation, visual_desc


class AudioExtractor(MediaExtractor):
    """Transcribe audio files using Gemini multimodal API."""

    @property
    def supported_mime_types(self) -> list[str]:
        return [
            "audio/mpeg",
            "audio/mp3",
            "audio/wav",
            "audio/x-wav",
            "audio/mp4",
            "audio/m4a",
            "audio/ogg",
            "audio/x-m4a",
        ]

    @property
    def supported_extensions(self) -> list[str]:
        return ["mp3", "wav", "m4a", "ogg", "flac", "aac"]

    async def extract(
        self, data: bytes, filename: str, metadata: dict[str, Any] | None = None
    ) -> MediaContent:
        settings = get_settings()
        size_mb = len(data) / (1024 * 1024)
        max_size = settings.media_max_file_size_mb

        if size_mb > max_size:
            return MediaContent(
                text=f"[Attachment: {filename} (audio, {size_mb:.1f} MB — exceeds {max_size} MB limit)]",
                media_type="audio",
            )

        parts: list[str] = [f"[Attachment: {filename} (audio, {size_mb:.1f} MB)]"]

        if not settings.google_api_key:
            return MediaContent(text="\n".join(parts), media_type="audio")

        uploaded = None
        client = await _get_gemini_client()
        try:
            from google.genai import types as genai_types

            ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else "mp3"
            mime_map = {
                "mp3": "audio/mpeg",
                "wav": "audio/wav",
                "m4a": "audio/mp4",
                "ogg": "audio/ogg",
                "flac": "audio/flac",
            }
            mime_type = mime_map.get(ext, "audio/mpeg")

            # Upload via Files API
            logger.info("AudioExtractor: uploading %s (%s, %.1f MB)", filename, mime_type, size_mb)
            uploaded = await asyncio.wait_for(
                client.aio.files.upload(
                    file=io.BytesIO(data),
                    config=genai_types.UploadFileConfig(mime_type=mime_type, display_name=filename),
                ),
                timeout=60,
            )
            file_uri: str = uploaded.uri or ""
            logger.info(
                "AudioExtractor: uploaded %s → %s (state=%s)",
                filename,
                uploaded.name,
                uploaded.state,
            )

            # Poll until file is ready for use
            await _poll_file_active(client, uploaded.name)
            logger.info("AudioExtractor: file ACTIVE for %s", filename)

            async with GEMINI_LIMITER:
                response = await asyncio.wait_for(
                    client.aio.models.generate_content(
                        model=settings.media_vision_model,
                        contents=[
                            genai_types.Content(
                                role="user",
                                parts=[
                                    genai_types.Part.from_uri(
                                        file_uri=file_uri,
                                        mime_type=mime_type,
                                    ),
                                    genai_types.Part.from_text(text=AUDIO_TRANSCRIPTION_PROMPT),
                                ],
                            )
                        ],
                    ),
                    timeout=60,
                )
            if response.text:
                logger.info(
                    "AudioExtractor: result for %s — %d chars", filename, len(response.text)
                )
                parts.append(f"[Audio summary]: {response.text}")
            else:
                logger.warning("AudioExtractor: empty response for %s", filename)
        except Exception:
            logger.error("AudioExtractor: transcription failed for %s", filename, exc_info=True)
        finally:
            if uploaded and getattr(uploaded, "name", None):
                try:
                    await client.aio.files.delete(name=uploaded.name)
                except Exception:
                    logger.debug("AudioExtractor: cleanup failed for %s", filename)

        return MediaContent(text="\n".join(parts), media_type="audio")

    @staticmethod
    def _parse_transcript(text: str) -> tuple[str, str]:
        """Parse structured transcript into (original, translation)."""
        transcript = ""
        translation = ""
        current_section = ""
        section_lines: dict[str, list[str]] = {"transcript": [], "translation": []}

        for line in text.split("\n"):
            stripped = line.strip()
            upper = stripped.upper()
            if upper.startswith("TRANSCRIPT:"):
                current_section = "transcript"
                rest = stripped[len("TRANSCRIPT:") :].strip()
                if rest:
                    section_lines["transcript"].append(rest)
            elif upper.startswith("TRANSLATION"):
                current_section = "translation"
                rest = stripped.split(":", 1)[-1].strip() if ":" in stripped else ""
                if rest:
                    section_lines["translation"].append(rest)
            elif current_section and stripped:
                section_lines[current_section].append(stripped)

        transcript = "\n".join(section_lines["transcript"]).strip()
        translation = "\n".join(section_lines["translation"]).strip()

        # Fallback: if no sections parsed, treat whole text as transcript
        if not transcript and text.strip():
            transcript = text.strip()

        return transcript, translation


def create_default_registry() -> MediaExtractorRegistry:
    """Create a registry with all default extractors registered."""
    registry = MediaExtractorRegistry()
    registry.register(PdfExtractor())
    registry.register(ImageExtractor())
    registry.register(OfficeExtractor())
    registry.register(VideoExtractor())
    registry.register(AudioExtractor())
    return registry
